# =====================================================
# 🤖 RAG Chatbot — Streamlit + Local Embeddings + Gemini
# User Login + Signup + Admin Dashboard + Usage Tracking
# =====================================================

import os
import json
import hashlib
import datetime
from dotenv import load_dotenv

load_dotenv()

import streamlit as st
import google.generativeai as genai
import numpy as np
import faiss
from PyPDF2 import PdfReader
from docx import Document
from pptx import Presentation
import pandas as pd
from langchain_huggingface import HuggingFaceEmbeddings

# =====================================================
# 🔐 CONFIG
# =====================================================

ADMIN_EMAIL = "Admin"
ADMIN_PASSWORD = "Poiu1234@@"
DB_FILE = "users_db.json"

# Gemini config
api_key = os.getenv("GOOGLE_API_KEY")
if not api_key:
    raise RuntimeError("❌ Missing GOOGLE_API_KEY in .env or Streamlit Secrets.")

genai.configure(api_key=api_key)
CHAT_MODEL = "gemini-2.0-flash"

# =====================================================
# 📁 USER DATABASE FUNCTIONS
# =====================================================

def init_db():
    if not os.path.exists(DB_FILE):
        with open(DB_FILE, "w") as f:
            json.dump({}, f)


def load_db():
    init_db()
    with open(DB_FILE, "r") as f:
        return json.load(f)


def save_db(db):
    with open(DB_FILE, "w") as f:
        json.dump(db, f, indent=4)


def hash_password(password):
    return hashlib.sha256(password.encode()).hexdigest()


def create_user(email, password):
    db = load_db()
    if email in db:
        return False  # already exists
    db[email] = {
        "password": hash_password(password),
        "login_count": 0,
        "queries": 0,
        "documents_uploaded": 0,
        "chunks_processed": 0,
        "last_login": "",
        "created": str(datetime.datetime.now())
    }
    save_db(db)
    return True


def authenticate_user(email, password):
    db = load_db()
    if email in db and db[email]["password"] == hash_password(password):
        return True
    return False


def update_usage(email, field, amount=1):
    db = load_db()
    if email in db:
        db[email][field] += amount
        save_db(db)


# =====================================================
# DOCUMENT HELPERS
# =====================================================

def extract_text(file):
    text = ""

    if file.name.endswith(".pdf"):
        reader = PdfReader(file)
        for page in reader.pages:
            text += page.extract_text() or ""

    elif file.name.endswith(".docx"):
        doc = Document(file)
        text = "\n".join([p.text for p in doc.paragraphs])

    elif file.name.endswith(".pptx"):
        prs = Presentation(file)
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text += shape.text + "\n"

    elif file.name.endswith(".txt"):
        text = file.read().decode("utf-8")

    elif file.name.endswith(".xlsx"):
        df = pd.read_excel(file)
        text = df.to_string()

    return text


def text_to_chunks(text, chunk_size=800, overlap=150):
    chunks = []
    start = 0
    while start < len(text):
        chunks.append(text[start:start+chunk_size])
        start += chunk_size - overlap
    return chunks


hf_embedder = HuggingFaceEmbeddings(
    model_name="sentence-transformers/all-MiniLM-L6-v2"
)

def embed_texts(texts):
    return np.array(hf_embedder.embed_documents(texts), dtype="float32")


def get_relevant_chunks(query, index, texts, k=3):
    qvec = np.array(hf_embedder.embed_query(query), dtype="float32").reshape(1, -1)
    _, idx = index.search(qvec, k)
    return [texts[i] for i in idx[0]]


# =====================================================
# LOGIN FLOWS
# =====================================================

def login_screen():
    st.title("🔐 User Login")

    email = st.text_input("Email")
    password = st.text_input("Password", type="password")

    if st.button("Login"):
        if email == ADMIN_EMAIL and password == ADMIN_PASSWORD:
            st.session_state["is_admin"] = True
            st.session_state["user"] = "admin"
            st.session_state["mode"] = "admin"
            st.rerun()

        if authenticate_user(email, password):
            st.session_state["user"] = email
            st.session_state["mode"] = "user"
            update_usage(email, "login_count")
            db = load_db()
            db[email]["last_login"] = str(datetime.datetime.now())
            save_db(db)
            st.success("Login successful!")
            st.rerun()
        else:
            st.error("Invalid login credentials.")

    if st.button("Create New Account"):
        st.session_state["mode"] = "signup"
        st.rerun()


def signup_screen():
    st.title("📝 Create Account")

    email = st.text_input("Email")
    password = st.text_input("Password", type="password")
    confirm = st.text_input("Confirm Password", type="password")

    if st.button("Create Account"):
        if password != confirm:
            st.error("Passwords do not match.")
        elif create_user(email, password):
            st.success("Account created! Please login now.")
            st.session_state["mode"] = "login"
            st.rerun()
        else:
            st.error("User already exists.")

    if st.button("⬅ Back to Login"):
        st.session_state["mode"] = "login"
        st.rerun()


def logout_button():
    if st.sidebar.button("🚪 Logout"):
        st.session_state["mode"] = "login"
        st.session_state["user"] = None
        st.session_state["is_admin"] = False
        st.rerun()


# =====================================================
# ADMIN PANEL
# =====================================================

def admin_panel():
    logout_button()

    st.title("🛠 Admin Dashboard")

    db = load_db()

    st.subheader("📊 Registered Users")
    st.table(db)

    if st.button("Clear All Users (Danger)"):
        save_db({})
        st.success("All users deleted.")

    if st.button("Clear Usage (Keep Accounts)"):
        db = load_db()
        for u in db:
            db[u]["queries"] = 0
            db[u]["documents_uploaded"] = 0
            db[u]["chunks_processed"] = 0
        save_db(db)
        st.success("Usage reset.")


# =====================================================
# USER APP (RAG CHATBOT)
# =====================================================

def user_app():

    logout_button()
    email = st.session_state["user"]

    st.title("🤖 Gemini RAG Chatbot")

    st.sidebar.header("📂 Upload Documents")
    files = st.sidebar.file_uploader(
        "Upload: PDF, DOCX, PPTX, TXT, XLSX",
        type=["pdf", "docx", "pptx", "txt", "xlsx"],
        accept_multiple_files=True
    )

    if files:
        update_usage(email, "documents_uploaded", len(files))

        with st.spinner("Processing documents..."):
            full_text = ""
            for f in files:
                full_text += extract_text(f) + "\n"

            chunks = text_to_chunks(full_text)
            update_usage(email, "chunks_processed", len(chunks))

            embeddings = embed_texts(chunks)
            index = faiss.IndexFlatL2(embeddings.shape[1])
            index.add(embeddings)

            st.session_state["faiss_index"] = index
            st.session_state["chunks"] = chunks

        st.success("Documents processed!")

    for role, msg in st.session_state.get("chat_history", []):
        st.markdown(f"**{role.upper()}:** {msg}")

    if "faiss_index" in st.session_state:
        query = st.text_input("Ask a question:")

        if query:
            update_usage(email, "queries")

            with st.spinner("Thinking..."):
                relevant = get_relevant_chunks(
                    query,
                    st.session_state["faiss_index"],
                    st.session_state["chunks"]
                )
                context = "\n".join(relevant)

                prompt = f"""
Answer ONLY using this context.
If answer not found: "I don’t have that information in the uploaded files."

Context:
{context}

Question:
{query}
"""

                model = genai.GenerativeModel(CHAT_MODEL)
                response = model.generate_content(prompt).text.strip()

                st.session_state.setdefault("chat_history", []).append(("user", query))
                st.session_state["chat_history"].append(("assistant", response))

                st.markdown(f"**Assistant:** {response}")

    else:
        st.info("Upload files to start chatting.")


# =====================================================
# APP ENTRY
# =====================================================

if "mode" not in st.session_state:
    st.session_state["mode"] = "login"

mode = st.session_state["mode"]

if mode == "login":
    login_screen()
elif mode == "signup":
    signup_screen()
elif mode == "admin":
    admin_panel()
elif mode == "user":
    user_app()
